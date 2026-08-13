#!/usr/bin/env python3
"""Build proposal/transakt-slides-4-6.pptx — slides 4, 5 and 6 of the Transakt deck.

    4  Product Architecture   <- proposal/product-architecture-slide.md
    5  Market                 <- proposal/market-slide.md
    6  Business Model Canvas  <- proposal/business-model-canvas.md

Those three markdown files are the source of truth. Edit them first, then mirror the
change here and rerun. Every figure traces to the provenance table at the foot of its
markdown file — do not change a number here without updating that table.

    pip install python-pptx
    python3 proposal/build_slides_4_5_6.py

Slide size is 13.333in x 7.5in, matching Pitch-Deck-June-8.pptx exactly, so these drop
straight into that deck in place of its slides 4, 5 and 6.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---- locked palette (identical across all three slides) ----------------------
NAVY   = RGBColor(0x0B, 0x1B, 0x3A)
TEAL   = RGBColor(0x0E, 0x8C, 0x7F)
BLUE   = RGBColor(0x1F, 0x6F, 0xB2)
AMBER  = RGBColor(0xC7, 0x7D, 0x18)
PURPLE = RGBColor(0x5A, 0x4F, 0xBF)
GREEN  = RGBColor(0x2C, 0x7A, 0x33)
GREY   = RGBColor(0x6B, 0x72, 0x80)
INK    = RGBColor(0x1B, 0x24, 0x30)
CREAM  = RGBColor(0xF7, 0xF4, 0xEF)
BORDER = RGBColor(0xEA, 0xE6, 0xDE)
HAIR   = RGBColor(0xD8, 0xD2, 0xC7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Aptos"

HERE = os.path.dirname(os.path.abspath(__file__))
OUT  = os.path.join(HERE, "transakt-slides-4-6.pptx")


# ---- primitives --------------------------------------------------------------
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, line=None, lw=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    return sh


def oval(slide, cx, cy, d, fill=None, line=None, lw=2.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2),
                                Inches(d), Inches(d))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    return sh


def text(slide, x, y, w, h, runs, size=9, color=INK, bold=False, italic=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, space_after=0):
    """runs: a string, or a list of paragraphs; each paragraph is a string or a list of
    (text, {overrides}) tuples."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        bits = para if isinstance(para, list) else [(para, {})]
        for t, ov in bits:
            r = p.add_run()
            r.text = t
            f = r.font
            f.name = FONT
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.italic = ov.get("italic", italic)
            f.color.rgb = ov.get("color", color)
    return tb


def bullets(slide, x, y, w, h, items, size=8, color=INK, spacing=1.15, gap=1.5):
    """items: list of strings, or (text, {overrides}). Rendered with a grey bullet glyph."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        t, ov = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = spacing
        p.space_after = Pt(gap)
        b = p.add_run()
        b.text = "•  "
        b.font.name, b.font.size, b.font.color.rgb = FONT, Pt(size), GREY
        r = p.add_run()
        r.text = t
        r.font.name = FONT
        r.font.size = Pt(ov.get("size", size))
        r.font.bold = ov.get("bold", False)
        r.font.color.rgb = ov.get("color", color)
    return tb


def hang(p, inches=0.10):
    """Hanging indent, so a wrapped bullet aligns under its text, not under the glyph."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(inches * 914400)))
    pPr.set("indent", str(-int(inches * 914400)))


def card(slide, x, y, w, h, accent, bar=3):
    """Cream card with a top accent bar — the deck's standard container."""
    rect(slide, x, y, w, h, fill=CREAM, line=BORDER, lw=0.75)
    rect(slide, x, y, w, bar / 72.0, fill=accent)


def kicker(slide, label, number):
    text(slide, 0.55, 0.34, 6.0, 0.28,
         [[(label, {"size": 11, "bold": True, "color": AMBER})]], size=11)
    text(slide, 12.0, 6.95, 0.78, 0.28, number, size=9, color=GREY, align=PP_ALIGN.RIGHT)


def title(slide, s):
    text(slide, 0.55, 0.72, 11.9, 0.62,
         [[(s, {"size": 26, "bold": True, "color": NAVY})]], size=26)


# ---- slide 4 — product architecture -----------------------------------------
def slide_04(prs):
    s = blank(prs)
    kicker(s, "PRODUCT ARCHITECTURE", "04")
    title(s, "Three trust tiers. Narrow on-chain, smart off-chain.")

    bands = [
        (1.60, 1.05, BLUE, "CLIENT TIER  ·  Buyer · Seller · Reviewer",
         "Every read and write goes through one authenticated API. Wallets sign deposit "
         "and release straight to the chain."),
        (2.95, 1.45, TEAL, "OFF-CHAIN ORCHESTRATION + VERIFICATION",
         "deal-intake · KYC/sanctions · document-checker + rules engine · dispute · settlement. "
         "AI reads, deterministic code computes the money, a human signs off. The append-only "
         "audit ledger is the regulator-facing source of truth."),
        (4.75, 1.05, NAVY, "ON-CHAIN TIER  ·  narrow Escrow contract  ·  Base Sepolia (EVM L2)",
         "Holds USDC · enforces the state machine · emits events. It never sees a trade document."),
    ]
    for y, h, accent, head, body in bands:
        card(s, 0.55, y, 8.10, h, accent)
        # one frame, vertically centred — keeps short and long bands equally balanced
        text(s, 0.67, y + 0.10, 7.86, h - 0.20,
             [[(head, {"size": 11, "bold": True, "color": accent})],
              [(body, {"size": 9, "color": INK})]],
             size=9, spacing=1.18, space_after=5, anchor=MSO_ANCHOR.MIDDLE)

    # the single releaser bridge — the slide's most important mark, sitting in the
    # gap between the off-chain band (ends 4.40) and the on-chain band (starts 4.75)
    ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.30), Inches(4.44),
                            Inches(0.30), Inches(0.27))
    ar.shadow.inherit = False
    ar.fill.solid()
    ar.fill.fore_color.rgb = AMBER
    ar.line.fill.background()
    text(s, 1.75, 4.49, 6.85, 0.20,
         [[("releaser key — the ONE bridge: off-chain verdict → recordVerdict on-chain",
            {"size": 8, "bold": True, "color": AMBER})]], size=8)

    # design-decision chips
    text(s, 8.95, 1.28, 3.85, 0.22,
         [[("DESIGN DECISIONS — WHY THE BUILD HOLDS", {"size": 8, "bold": True, "color": NAVY})]],
         size=8)
    chips = [
        ("Narrow contract", "AP-1",
         "On-chain code holds funds and state only — no document logic, no money math. "
         "Small attack surface, chain-portable."),
        ("Permissionless release", "AP-7",
         "Release isn't gated on our key — a cleared seller can always be paid. "
         "Liveness without trusting the operator."),
        ("Audit around every action", "AP-4",
         "Intent logged before the transaction, reconciliation after. The ledger — not the "
         "chain — is the source of truth."),
        ("Money in code, human above the line", "AP-5 / 7",
         "All arithmetic is deterministic code; anything outside the autonomy thresholds "
         "escalates to a human sign-off."),
    ]
    cy = 1.60
    for head, code, body in chips:
        rect(s, 8.95, cy, 3.85, 1.05, fill=CREAM, line=BORDER, lw=0.75)
        rect(s, 8.95, cy, 2 / 72.0, 1.05, fill=NAVY)
        text(s, 9.10, cy + 0.11, 3.60, 0.20,
             [[(head + "  ", {"size": 9.5, "bold": True, "color": NAVY}),
               (code, {"size": 8, "bold": True, "color": AMBER})]], size=9.5)
        text(s, 9.10, cy + 0.34, 3.60, 0.62, body, size=8, color=INK, spacing=1.15)
        cy += 1.17

    text(s, 0.55, 6.28, 8.10, 0.30,
         [[("Narrow on-chain, smart off-chain, audited around every action — "
            "documentary-grade trust without the issuing bank.",
            {"size": 11, "italic": True, "color": NAVY})]], size=11)
    return s


# ---- slide 5 — market --------------------------------------------------------
def slide_05(prs):
    s = blank(prs)
    kicker(s, "MARKET", "05")
    title(s, "The wedge is narrow. The trade behind it is not.")

    # centre raised to 3.68 so the outer ring (Ø4.05, bottom 5.71) clears the
    # supporting line at 6.02 — at Ø4.30 centred on 4.05 the two collided
    cx, cy = 4.05, 3.68
    oval(s, cx, cy, 4.05, fill=None, line=NAVY, lw=2.5)
    oval(s, cx, cy, 2.70, fill=None, line=BLUE, lw=2.5)
    oval(s, cx, cy, 1.40, fill=TEAL, line=None)

    # TAM / SAM labels on the upper arc of their own ring; SOM inside the disc
    text(s, cx - 1.85, cy - 1.90, 3.70, 0.46,
         [[("$1.03T", {"size": 20, "bold": True, "color": NAVY})],
          [("UAE non-oil foreign trade, 2025", {"size": 8.5, "color": GREY})]],
         align=PP_ALIGN.CENTER, spacing=1.0)
    text(s, cx - 1.60, cy - 1.22, 3.20, 0.46,
         [[("~$65B", {"size": 20, "bold": True, "color": BLUE})],
          [("UAE–India CEPA non-oil corridor", {"size": 8.5, "color": GREY})]],
         align=PP_ALIGN.CENTER, spacing=1.0)
    text(s, cx - 0.68, cy - 0.28, 1.36, 0.60,
         [[("~$110M", {"size": 15, "bold": True, "color": WHITE})],
          [("processed, Y1–Y3", {"size": 7.5, "color": WHITE})]],
         align=PP_ALIGN.CENTER, spacing=1.0)

    # tier keys, left of the rings, each level with its own label row
    for ty, lbl, col in ((cy - 1.86, "TAM", NAVY), (cy - 1.18, "SAM", BLUE),
                         (cy - 0.24, "SOM", TEAL)):
        text(s, 0.60, ty, 0.80, 0.24,
             [[(lbl, {"size": 11, "bold": True, "color": col})]], size=11)

    # beachhead card — vertical centre matched to the circles' centre (3.68)
    card(s, 7.30, 2.48, 5.45, 2.40, AMBER)
    text(s, 7.46, 2.66, 5.13, 0.24,
         [[("BEACHHEAD LOGIC", {"size": 9, "bold": True, "color": AMBER})]], size=9)
    lines = [
        ("1", "UAE importers of Indian textiles, garments and consumer goods."),
        ("2", "$10K–$50K shipments — the LC dead zone: too small for a bank to issue "
              "against, too large to pay on trust."),
        ("3", "Repeat trade pattern creates a retention surface."),
    ]
    ly = 3.00
    for num, body in lines:
        text(s, 7.46, ly, 0.30, 0.30,
             [[(num, {"size": 13, "bold": True, "color": AMBER})]], size=13)
        text(s, 7.84, ly + 0.03, 4.75, 0.62, body, size=10.5, color=INK, spacing=1.18)
        ly += 0.72

    text(s, 0.55, 6.02, 6.95, 0.28,
         [[("Where the gap bites hardest, and where the regulators are most ready.",
            {"size": 12, "italic": True, "color": NAVY})]], size=12, align=PP_ALIGN.CENTER)
    text(s, 0.55, 6.42, 6.95, 0.24,
         [[("Trade value we settle against — not a financing gap we fund.",
            {"size": 8.5, "italic": True, "color": GREY})]], size=8.5, align=PP_ALIGN.CENTER)
    text(s, 0.55, 7.02, 9.5, 0.22,
         [[("Sources: UAE Ministry of Foreign Trade (2026); India–UAE CEPA bilateral trade "
            "(2026); Blockmediary financial model, 13 Aug 2026. SOM converted at a stated "
            "1.30 USD/GBP.", {"size": 6.5, "italic": True, "color": GREY})]], size=6.5)
    return s


# ---- slide 6 — business model canvas ----------------------------------------
def slide_06(prs):
    s = blank(prs)
    text(s, 0.30, 0.34, 7.0, 0.40,
         [[("Business Model Canvas", {"size": 18, "bold": True, "color": NAVY})]], size=18)
    text(s, 6.0, 0.44, 7.03, 0.26,
         [[("Transakt — Blockmediary", {"size": 10, "italic": True, "color": GREY})]],
         size=10, align=PP_ALIGN.RIGHT)
    rect(s, 0.30, 0.90, 12.73, 1 / 72.0, fill=HAIR)

    BW, BH, GAP = 3.11, 2.90, 0.10
    cols = [0.30 + i * (BW + GAP) for i in range(4)]
    rows = [1.00, 4.00]

    def block(col, row, accent, head, question, sections):
        """One card, one text frame. Every line is a paragraph in a single flowing
        frame, so PowerPoint does the wrapping — no hand-estimated line heights."""
        x, y = cols[col], rows[row]
        card(s, x, y, BW, BH, accent)

        tb = s.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.14),
                                  Inches(BW - 0.24), Inches(BH - 0.26))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        paras = []          # (runs, line_spacing, space_before, space_after)
        paras.append(([(head, 11, True, False, NAVY)], 1.10, 0, 1.5))
        if question:
            paras.append(([(question, 7, False, True, GREY)], 1.08, 0, 3.0))
        for label, items in sections:
            if label:
                paras.append(([(label, 8, True, False, accent)], 1.10, 3.0, 1.0))
            for it in items:
                t, ov = it if isinstance(it, tuple) else (it, {})
                paras.append(([("•  ", 8, False, False, GREY),
                               (t, 8, ov.get("bold", False), False, INK)], 1.12, 0, 1.2))

        for i, (runs, ls, sb, sa) in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = ls
            p.space_before = Pt(sb)
            p.space_after = Pt(sa)
            if runs[0][0] == "•  ":
                hang(p)
            for t, sz, bold, ital, col_ in runs:
                r = p.add_run()
                r.text = t
                r.font.name = FONT
                r.font.size = Pt(sz)
                r.font.bold = bold
                r.font.italic = ital
                r.font.color.rgb = col_

    block(0, 0, NAVY, "The Business Opportunity",
          "Where is the gap in the market, and how do you intend to fill it?",
          [("The Problem", [
              ("Same fixed compliance cost on a $10K trade as a $5M trade", {"bold": True}),
              "SME ticket sizes uneconomic for banks at scale",
              "$2.5T global trade finance gap (ADB, 2025); 41% of SME applications rejected",
              "Rejected SMEs fall back on trust, prepayment, or drop the trade"]),
           ("Our Solution", [
              "UCP 600-inspired documentary release, stablecoin settlement, no issuing bank",
              "Buyer escrow → seller documents → release on compliant presentation",
              "Four verification layers plus an invoked escalation route",
              "Blockmediary addresses the documentary trust and payment-assurance gap, "
              "not working-capital finance"])])

    block(1, 0, PURPLE, "Business Description", "Brief outline of your business model",
          [(None, [
              ("Documentary escrow operator: verification, release rules, settlement", {"bold": True}),
              "Smart-contract escrow holds stablecoins against documentary release rules",
              "UCP 600-inspired Article 14 release logic — a Trade Escrow Agreement, not an LC",
              "Settlement on stablecoin rails (USDC), not bank correspondent networks",
              "Not a bank: no deposits, no lending, no token issuance",
              "Not DeFi: centralised trust and compliance layer by design",
              "Pilot planned through a suitably DFSA-authorised partner, subject to confirmed "
              "permissions and executed agreements",
              "First paid deal targeted Month 12; full VARA licence Month 18 — the gate for "
              "UAE-wide scale"])])

    block(2, 0, BLUE, "Target Market", None,
          [("Sector & geography", [
              ("SME importers/exporters; UAE launch market", {"bold": True}),
              "Corridor: UAE–India (CEPA); more added per pilot envelope",
              "Sectors: manufacturing, consumer goods, electronics, textiles",
              "Ticket: the $10K–$50K “LC dead zone”; average deal £31K → £48K by Y5"]),
           ("Strategic goals", [
              "Payment without LC cost or bank issuance delay",
              "Counterparty access without bank gatekeeping",
              "Verified trade reputation that travels across deals"]),
           ("Pain points & risks", [
              "LC pricing and access barrier at SME tickets",
              "Pre-shipment payment-trust gap",
              "Document fraud and discrepancy risk",
              "Corridor FX and regulatory friction"])])

    block(3, 0, TEAL, "USPs",
          "How does your product/service compare? What makes you unique?",
          [(None, [
              ("UCP 600-inspired documentary release on stablecoin rails", {"bold": True}),
              "Four-layer verification: screening → extraction → source corroboration "
              "→ contracted examiner",
              "Customer-selected verification tiers — speed traded against cost of trust",
              "Two-track UAE route: DIFC partner pilot, then own VARA licence for scale",
              "Modelled to undercut traditional LC fees on Tier A and selected Tier B transactions",
              "Compliance-first by design: not a bank, not DeFi"])])

    block(0, 1, GREEN, "Revenue", None,
          [("Streams", [
              ("A. Tiered escrow fee: 0.8% / 1.5% / 3.0% by verification tier", {"bold": True}),
              "B. Document review, dispute, expedited and onboarding fees",
              "C. Partner/API setup and committed-volume fees"]),
           ("Projected revenue — base case", [
              "Year 3: £1.37m revenue on 2,000 deals and £76m GMV",
              "Blended take rate falls 2.0% → 1.0% as Tier A grows 15% → 70%",
              "Gross margin 71–79%; ancillary fees 25% of Y3 revenue, 17% by Y5",
              "Monthly cash-flow break-even in Year 4, Month 8"])])

    block(1, 1, AMBER, "Costs", None,
          [("Funding requirement", [
              ("£3.06m operating funding to break-even; £319k VARA capital separate",
               {"bold": True})]),
           ("Fixed", [
              "Payroll — 8 FTE at launch to 24 by Year 5; dominant cost",
              "Compliance, MLRO, VARA supervision and operational resilience",
              "Engineering, cloud, monitoring and security tooling",
              "Insurance (PI, cyber, D&O) and legal retainer"]),
           ("Variable (per deal)", [
              "eBL/carrier API lookup, or paper collection",
              "Contracted documentary examiner time",
              "KYB/KYC screening per new counterparty",
              "On-chain gas and settlement cost",
              "Dispute handling and enhanced review above £50K"])])

    block(2, 1, BLUE, "Channels", "How will you promote your product or service?",
          [("Direct", [
              ("Dubai and Sharjah chambers of commerce", {"bold": True}),
              "Freight forwarder referrals",
              "DIFC FinTech Hive; Gulfood, GITEX, Dubai FinTech Summit"]),
           ("Indirect", [
              "Regulated-partner distribution — the base case's primary scale route",
              "Fintech wallet integrations (white-label verification API)",
              "Logistics platform API partnerships"])])

    block(3, 1, NAVY, "Competition", "Who are your competitors and what are their USPs?",
          [(None, [
              "Bank-issued LCs — trust and networks; 1–3% commission before the fee stack",
              "Tazapay — escrow-as-a-service; manual review, custodial fiat rails",
              "Truzo — narrow UK–South Africa corridor; escrow wallet, no document layer",
              "XREX — licensed custodial escrow; releases on agreement, not documents",
              "Komgo / Contour — bank-consortium LC digitisation; Contour shut in 2023"])])

    text(s, 0.30, 6.94, 12.73, 0.34,
         [[("Sources: ADB (2025); ICC UCP 600 (2007); VARA Rulebook (2026); Blockmediary "
            "financial model and legal & compliance risk register, 13 Aug 2026. UCP 600-inspired "
            "release logic; not a letter of credit. Pilot planned through a suitably "
            "DFSA-authorised partner, subject to confirmed permissions and executed agreements. "
            "Blockmediary provides documentary escrow and payment assurance, not working-capital "
            "finance. Figures are base case. GMV is trade value processed, not revenue.",
            {"size": 6, "italic": True, "color": GREY})]], size=6, spacing=1.35)
    return s


def main():
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.333in
    prs.slide_height = Emu(6858000)   # 7.5in
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    prs.save(OUT)
    print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
